# app.py
import os, os.path
import re
import uuid
from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify
from werkzeug.utils import secure_filename
from pypdf import PdfReader
import docx
import shutil
import stat
from pptx import Presentation
import time
import random
import io
import logging
from threading import Thread
from datetime import datetime
from tempfile import SpooledTemporaryFile
from flask_mail import Mail, Message
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import CharacterTextSplitter
from openai import OpenAI
from dotenv import load_dotenv


# load environment variables from .env for the OpenAI code
#load_dotenv(os.path.join(os.path.dirname(__file__), ".env"))

# ---------- Flask setup ----------
BASEDIR = os.path.dirname(__file__)
load_dotenv(os.path.join(BASEDIR, ".env"))  # load .env before creating the app

app = Flask(__name__)

# ✅ REQUIRED for sessions/flash — try FLASK_SECRET_KEY, then SECRET_KEY, then a dev fallback
app.config["SECRET_KEY"] = (
    os.getenv("FLASK_SECRET_KEY")
    or os.getenv("SECRET_KEY")
    or "dev-change-me"
)

# Flask-Mail config
app.config.update(
    MAIL_SERVER="smtp.gmail.com",
    MAIL_PORT=587,
    MAIL_USE_TLS=True,

    MAIL_USERNAME=os.getenv("MAIL_USERNAME"),
    MAIL_PASSWORD=os.getenv("MAIL_PASSWORD"),
    MAIL_DEFAULT_SENDER=os.getenv("MAIL_USERNAME")
)

mail = Mail(app)


app.config["UPLOAD_FOLDER"] = os.path.join(BASEDIR, "uploads")
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

# Allowed extensions (unchanged)
ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx"}


# ---------- Helpers ----------

#helpers for logging
log = logging.getLogger("timing")
log.setLevel(logging.INFO)
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s"
)

PROGRESS = {} 

def set_progress(job_id, phase, pct):
    PROGRESS[job_id] = {**PROGRESS.get(job_id, {}), "phase": phase, "pct": int(pct)}
    logging.info("SET_PROGRESS job=%s phase=%s pct=%s", job_id, phase, pct)


class Timer:
    def __init__(self, name):
        self.name = name
        self.start = time.perf_counter()

    def done(self, extra=""):
        elapsed = time.perf_counter() - self.start
        log.info("[TIMER] %-20s %7.3fs %s", self.name, elapsed, extra)
        return elapsed


def allowed_file(filename: str) -> bool:
    """Check extension against the allowed set."""
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF pages (no OCR)."""
    try:
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            return "".join([(page.extract_text() or "") for page in reader.pages])
    except Exception as e:
        return f"Error extracting text from PDF: {e}"


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX paragraphs."""
    try:
        d = docx.Document(file_path)
        return "\n".join([p.text for p in d.paragraphs if p.text.strip()])
    except Exception as e:
        return f"Error extracting text from DOCX: {e}"

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    return "\n".join(parts).strip()


def process_uploaded_file(file_storage):
    """
    Save the uploaded file securely and extract its text content.
    Returns (text, base_name_without_ext, sanitized_filename).
    """
    filename = secure_filename(file_storage.filename)
    save_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file_storage.save(save_path)

    lower = filename.lower()
    if lower.endswith(".pdf"):
        text = extract_text_from_pdf(save_path)
    elif lower.endswith(".docx"):
        text = extract_text_from_docx(save_path)
    elif lower.endswith(".pptx"):
        text = extract_text_from_pptx(save_path)
    elif lower.endswith(".txt"):
        with open(save_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    else:
        text = "Unsupported file type."

    base_name_without_ext = os.path.splitext(filename)[0]
    return text, base_name_without_ext, filename


def get_openai_client():
    """Create an OpenAI client using the API key in the environment."""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY not set in environment.")
    return OpenAI(api_key=api_key)


def get_document_prompt(docs):
    """Format a list of strings or LangChain Documents into a numbered prompt block."""
    out = []
    for i, d in enumerate(docs, 1):
        text = d if isinstance(d, str) else getattr(d, "page_content", "")
        out.append(f"\nContent {i}:\n{text}\n")
    return "\n".join(out)



def _on_rm_error(func, path, exc_info):
    """Windows-safe remover: make file writable then retry."""
    try:
        os.chmod(path, stat.S_IWRITE)
        func(path)
    except Exception:
        logging.warning("Could not remove: %s", path)


def _process_job(job_id: str, text: str, persist_dir: str, filename: str, start_pct: int = 40, end_pct: int = 100):
    """
    Run the embedding + summary build and report progress scaled into [start_pct, end_pct].
    """
    def scale(local):  # local is 0..100 → map into [start..end]
        local = max(0, min(100, int(local)))
        span = max(1, end_pct - start_pct)
        return start_pct + int(round(local * span / 100.0))

    try:
        PROGRESS[job_id] = {"phase": "Processing", "pct": scale(2)}
        os.makedirs(persist_dir, exist_ok=True)

        # Split text
        PROGRESS[job_id] = {"phase": "Processing", "pct": scale(5)}
        text_splitter = CharacterTextSplitter(separator=" ", chunk_size=5000, chunk_overlap=100)
        docs = text_splitter.split_text(text) if text else []

        # Build embeddings + DB
        PROGRESS[job_id] = {"phase": "Processing", "pct": scale(10)}
        client = get_openai_client()
        embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=client.api_key)
        vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)

        total = max(len(docs), 1)
        batch = 50
        added = 0
        for i in range(0, len(docs), batch):
            chunk = docs[i:i + batch]
            vectordb.add_texts(chunk)
            added += len(chunk)
            local_pct = 10 + (65 * added / total)  # 10→75 locally
            PROGRESS[job_id] = {"phase": f"Processing", "pct": scale(local_pct)}

        PROGRESS[job_id] = {"phase": "Processing", "pct": scale(20)}

        PROGRESS[job_id] = {"phase": "Processing", "pct": scale(30)}


        # Summarize
        PROGRESS[job_id] = {"phase": "Summarizing", "pct": scale(90)}
        raw = vectordb.get(include=["documents"])
        sample = (raw.get("documents") or [])[:15]
        prompt = get_document_prompt(sample) if sample else "No content available."

        system_message = (
            f"Generate a summary of the following notebook content::\n\n"
            f"\n\n###\n{prompt}\n###\n\n"
            "The summary should contain the title of the book and a short sentence about the notebook"
            "The first line must be the notebook title, wrapped in double asterisks, like:\n"
            "Title: at the beginning of the notebook title"
            "Then add two newline characters (\\n\\n).\n"
            "After that, write the rest of the summary"
            "The summary should never be more that 2 sentences"
            "Be precise, avoid opinions, and summarize the main points in a clear and structured way. "
            "If the document has multiple sections, break it into meaningful segments."
        )
        resp = get_openai_client().chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": system_message}],
            temperature=0.2,
        )
        summary_text = resp.choices[0].message.content

        PROGRESS[job_id] = {"phase": "completed", "pct": scale(100), "summary": summary_text, "filename": filename, }

    except Exception as e:
        PROGRESS[job_id] = {"phase": "error", "pct": end_pct, "error": str(e)}




# ---------- Routes ----------
@app.get("/")
@app.get("/home")
def home():
    #landing page for uploading document
    return render_template("home.html")


@app.post("/delete_doc")
def delete_doc():
    # Accept JSON (from fetch) and also allow form (fallback)
    data = request.get_json(silent=True) or {}
    filename = data.get("filename") or request.form.get("filename")

    if not filename:
        return jsonify(ok=False, error="Please select a document to delete."), 400

    docs = session.get("docs", {})

    if filename not in docs:
        return jsonify(ok=False, error="Selected document not found."), 404

    persist_dir = docs[filename].get("persist_dir")
    if not persist_dir:
        return jsonify(ok=False, error="No database path stored for this document."), 500
    
    #important on Windows
    persist_dir = os.path.abspath(persist_dir)  
    app.logger.info("DELETE requested for %s, persist_dir=%s", filename, persist_dir)

    if not os.path.isdir(persist_dir):
        # Don’t pretend it worked
        return jsonify(ok=False, error=f"Database folder not found: {persist_dir}"), 404

    # Try deleting with retries (Windows file locks)
    last_err = None
    for attempt in range(3):
        try:
            shutil.rmtree(persist_dir, onerror=_on_rm_error)
            last_err = None
            break
        except Exception as e:
            last_err = e
            app.logger.warning("rmtree attempt %s failed: %s", attempt + 1, e)
            time.sleep(0.3)

    if last_err:
        return jsonify(ok=False, error=f"Failed to delete database: {last_err}"), 500

    # Remove from session
    docs.pop(filename, None)
    session["docs"] = docs

    # Clear active doc if it was the deleted one
    if session.get("uploaded_filename") == filename:
        session.pop("uploaded_filename", None)
        session.pop("summary_text", None)
        session.pop("persist_directory", None)

    flash(f"Deleted '{filename}' successfully.", "success")
    return jsonify(ok=True, message=f"Deleted '{filename}' successfully.")


#generate questions
@app.get("/upload_notebook")
def upload_notebook():
    job_id = session.get("job_id")
    if job_id:
        st = PROGRESS.get(job_id, {})
        phase = (st.get("phase") or "").lower()
        if phase == "completed" and st.get("summary"):
            docs = session.get("docs", {})
            filename = st.get("filename") or session.get("uploaded_filename")
            # Save summary per document
            if filename:
                info = docs.get(filename, {})
                info["persist_dir"] = info.get("persist_dir") or session.get("persist_directory")
                info["summary"] = st["summary"]
                docs[filename] = info
                session["docs"] = docs
                # Also keep "current" summary for the page
                session["uploaded_filename"] = filename
                session["summary_text"] = st["summary"]
            PROGRESS.pop(job_id, None)
            session.pop("job_id", None)
            job_id = None  # <- ensures the template won’t emit data-job-id
        elif phase == "error":
            flash(f"⚠️ Could not build index or generate summary: {st.get('error')}", "error")
            PROGRESS.pop(job_id, None)
            session.pop("job_id", None)
            job_id = None

    return render_template("upload_notebook.html",
        filename=session.get("uploaded_filename"),
        summary=session.get("summary_text"),
        job_id=job_id,  # will be None if finished
    )


#This is for clicking a checkbox and instantly seeing that file’s summary.
@app.get("/summary")
def get_summary():
    filename = request.args.get("filename", "")
    docs = session.get("docs", {})
    info = docs.get(filename) or {}
    summary = info.get("summary")
    return jsonify({"ok": True, "summary": summary or ""})


#Used to syle the summary
@app.template_filter("markdown_bold")
def markdown_bold(s):
    if not isinstance(s, str):
        return s
    return re.sub(r"\*\*(.*?)\*\*", r"<strong>\1</strong>", s)


@app.template_filter("nl2br")
def nl2br(s):
    if not isinstance(s, str):
        return s
    # turn newlines into <br> for HTML
    return s.replace("\n", "<br>")


#Create a job_id before uploading so the UI can start polling#
@app.post("/init_upload")
def init_upload():
    job_id = uuid.uuid4().hex
    PROGRESS[job_id] = {"phase": "Uploading", "pct": 0}
    session["job_id"] = job_id  # so /generate can render it
    return jsonify({"ok": True, "job_id": job_id})


@app.post("/upload")
def upload():
    job_id = request.headers.get("X-Job-Id") or session.get("job_id")
    if not job_id:
        job_id = uuid.uuid4().hex
    logging.info(
        "UPLOAD START job=%s header_job=%s content_length=%s",
        job_id,
        request.headers.get("X-Job-Id"),
        request.content_length,
    )

    session["job_id"] = job_id
    PROGRESS.setdefault(job_id, {"phase": "Uploading", "pct": 0})

    # Detect XHR upload (AJAX)
    is_xhr = request.headers.get("X-Requested-With") == "XMLHttpRequest"

    # Validation
    if "file" not in request.files:
        msg = "No file part in request."
        if is_xhr:
            return jsonify(ok=False, error=msg), 400
        flash(msg, "error")
        return redirect(url_for("upload_notebook"))

    f = request.files["file"]

    if f.filename == "":
        msg = "No file selected."
        if is_xhr:
            return jsonify(ok=False, error=msg), 400
        flash(msg, "error")
        return redirect(url_for("upload_notebook"))

    if not allowed_file(f.filename):
        msg = "Unsupported file type. Please upload PDF, DOCX, PPTX or TXT."
        if is_xhr:
            return jsonify(ok=False, error=msg), 400
        flash(msg, "error")
        return redirect(url_for("upload_notebook"))

    # Save + extract text
    text, base, filename = process_uploaded_file(f)

    # Here upload is complete → move to 40%
    PROGRESS[job_id] = {
        "phase": "queued",
        "pct": 40,
        "filename": filename
    }
    # Track uploaded list for sidebar
    uploaded = session.get("uploaded_files", [])
    uploaded.append(filename)
    session["uploaded_files"] = uploaded

    # Init session state for this upload
    persist_dir = os.path.abspath(
        f"./chroma_db_{base}_{uuid.uuid4().hex[:8]}"
    )

    docs = session.get("docs", {})
    docs[filename] = {
        "persist_dir": persist_dir,
    }
    session["docs"] = docs
    session["persist_directory"] = persist_dir
    session["uploaded_filename"] = filename
    session["summary_text"] = None
    session["summary_generated"] = False

    # Start background processing
    PROGRESS[job_id] = {
        "phase": "queued",
        "pct": 40,
        "filename": filename
    }

    t = Thread(
        target=_process_job,
        args=(job_id, text, persist_dir, filename, 40, 100),
        daemon=True
    )
    t.start()

    app.logger.info(
        "XHR=%s X-Job-Id=%s",
        is_xhr,
        request.headers.get("X-Job-Id")
    )
    # Response
    if is_xhr:
        return jsonify(
            ok=True,
            job_id=job_id,
            filename=filename
        ), 200

    return redirect(url_for("upload_notebook"))


@app.get("/progress/<job_id>")
def get_progress(job_id):
    st = PROGRESS.get(job_id, {"phase": "queued", "pct": 0})

    # Never allow pct to go backwards (safety)
    prev = PROGRESS.get(job_id, {})
    if "pct" in prev and "pct" in st:
        st["pct"] = max(int(prev.get("pct", 0)), int(st.get("pct", 0)))

    resp = jsonify(st)
    resp.headers["Cache-Control"] = "no-store"
    return resp





@app.route("/ask", methods=["POST"])
def ask():
    data = request.get_json(silent=True) or {}
    question = data.get("question", "").strip()
    filename = data.get("filename")
    if not question:
        return jsonify({"ok": False, "error": "Question is required."}), 400
    
    # Look up persist_dir by filename
    docs = session.get("docs", {})
    info = docs.get(filename or "", {})

    persist_dir = info.get("persist_dir") if info else None
    if not persist_dir or not os.path.isdir(persist_dir):
        return jsonify({"ok": False, "error": "Please select a Notebook before asking a Question."}), 400

    client = get_openai_client()
    embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=client.api_key)
    vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)

    # Retrieve relevant docs
    retrieved = vectordb.similarity_search(question, k=10)
    context = get_document_prompt(retrieved)

    system_message = (
        f"You are a professor teaching a course. Use the following notebook content "
        f"to answer student questions accurately and concisely:\n\n{context}\n\n"
        "Be precise and avoid opinions."
        "Only state what is in the notebook content"
        "Do not state what is not in the given notebook and be very precise and straight forward "
    )

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_message},
            {"role": "user", "content": question},
        ],
        temperature=0.1,
    )
    answer = resp.choices[0].message.content
    return jsonify({"ok": True, "answer": answer})


#Generate multiple-choice questions from the vector DB.
@app.route("/generate_quiz", methods=["POST"])
def generate_quiz():
    t_total = Timer("generate_quiz TOTAL")

    # ----------------------------
    t_request = Timer("parse request")
    data = request.get_json(silent=True) or {}
    num = int(data.get("num_questions", 5))
    filename = data.get("filename")
    t_request.done(f"(num={num})")

    # ----------------------------
    t_session = Timer("session lookup")
    # Look up persist_dir by filename
    docs = session.get("docs", {})
    info = docs.get(filename or "", {})
    persist_dir = info.get("persist_dir") if info else None
    t_session.done()

    if not persist_dir or not os.path.isdir(persist_dir):
        return jsonify({"ok": False, "error": "Please select a Notebook before generating Quiz."}), 400

    # ----------------------------
    t_vectordb = Timer("load vector DB")
    client = get_openai_client()
    embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=client.api_key)
    vectordb = Chroma(embedding_function=embeddings, persist_directory=persist_dir)
    t_vectordb.done()
    
    # ----------------------------
    t_fetch = Timer("fetch documents")
    raw = vectordb.get(include=["documents"])
    # Use at least 20 documents selection when available
    all_docs = raw.get("documents", [])
    t_fetch.done(f"(docs={len(all_docs)})")

    # ----------------------------
    t_sample = Timer("sample documents")
    num_samples = min(20, len(all_docs)) if all_docs else 0
    sample = random.sample(all_docs, num_samples) if num_samples > 0 else []
    context = get_document_prompt(sample) if sample else "No content available."
    t_sample.done(f"(sampled={num_samples}, chars={len(context)})")

    # ----------------------------
    t_prompt = Timer("build prompt")
    system_message = (
        f"Generate {num} multiple-choice quiz questions from the following notebook content: "
        f"\n\n###\n{context}\n###\n\n"
        f"Each question should have 4 answer choices (A,B,C,D) and indicate the correct answer at the end:"
        f"""the format of the reply should be
        Question 1: <question>
        A)  <answer choice A>
        B)  <answer choice B>
        C)  <answer choice C>
        D)  <answer choice D>
        Correct Answer: C

        Question 2: <question>
         ..."""
         )
    t_prompt.done()

    # ----------------------------
    t_llm = Timer("LLM generation")
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": system_message}],
        temperature=0.2,
    )
    t_llm.done()

    # ----------------------------
    t_parse = Timer("parse LLM output")
    text = resp.choices[0].message.content.strip()
    # Parse very simply
    blocks = [b for b in text.split("\n\n") if b.strip()]
    quiz = []
    for b in blocks:
        lines = b.splitlines()
        if len(lines) >= 6 and lines[0].lower().startswith("question"):
            q = lines[0].strip()
            choices = lines[1:5]
            correct = lines[5].split(":")[-1].strip()
            quiz.append({"question": q, "choices": choices, "correct": correct})
    t_parse.done(f"(parsed={len(quiz)})")

    # ----------------------------
    t_total.done()
    return jsonify({"ok": True, "quiz": quiz})

#for saving results
@app.post("/save_result")
def save_result():
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")
    correct = data.get("correct")
    total = data.get("total")
    percent = data.get("percent")

    if not filename or correct is None or total is None:
        return jsonify(ok=False, error="Missing result data"), 400

    results = session.get("results", [])
    results.insert(0, {
        "filename": filename,
        "correct": int(correct),
        "total": int(total),
        "percent": int(percent) if percent is not None else round((int(correct)/int(total))*100),
        "test_datetime": datetime.now().strftime("%Y-%m-%d %I:%M %p")
    })
    session["results"] = results
    return jsonify(ok=True)


#results page route
@app.get("/results")
def results():
    return render_template(
        "results.html",
        results=session.get("results", [])
    )

@app.route("/send-feedback", methods=["POST"])
def send_feedback():
    rating = request.form.get("rating", "N/A")
    category = request.form.get("category", "N/A")
    message = request.form.get("message", "").strip()

    if not message:
        return jsonify({
            "status": "error",
            "message": "Message cannot be empty."
        })

    try:
        msg = Message(
            subject=f"Studyassists Feedback — {category}",
            recipients=["info@studyassists.com"], 
            body=(
                f"Rating: {rating}\n"
                f"Category: {category}\n\n"
                f"Message:\n{message}"
            )
        )

        mail.send(msg)

        return jsonify({
            "status": "success",
            "message": "Thank you! Your feedback has been sent."
        })

    except Exception as e:
        print("Email error:", e)
        return jsonify({
            "status": "error",
            "message": "Unable to send feedback at this time."
        })





if __name__ == "__main__":
    import os
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, threaded=True)



